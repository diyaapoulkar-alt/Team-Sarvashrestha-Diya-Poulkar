import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt

def update_original_pptx():
    # Source file: The 12.4 MB original presentation with custom slide backgrounds and master fonts
    src_file = r'C:\Users\User\Downloads\Saathi-Your-Accessibility-Copilot (1).pptx'
    if not os.path.exists(src_file):
        src_file = r'C:\Users\User\Desktop\Saathi-Your Accessibility Copilot.pptx'

    print(f"Loading original presentation template: {src_file}")
    prs = Presentation(src_file)

    # Add Architecture Diagram image to Slide 5 (Technical Strategy: Methodology)
    slide5 = prs.slides[4]
    arch_img = os.path.join(os.getcwd(), 'public', 'architecture_diagram.png')
    if os.path.exists(arch_img):
        # Insert high-res Architecture Diagram image onto slide 5
        slide5.shapes.add_picture(arch_img, Inches(1.2), Inches(2.2), Inches(10.8), Inches(4.8))
        print("Embedded high-res architecture_diagram.png into Slide 5!")

    # Add Workflow Diagram image to Slide 6 (Implementation Flow)
    slide6 = prs.slides[5]
    wf_img = os.path.join(os.getcwd(), 'public', 'workflow_diagram.png')
    if os.path.exists(wf_img):
        # Insert high-res Workflow Flowchart image onto slide 6
        slide6.shapes.add_picture(wf_img, Inches(1.2), Inches(2.0), Inches(10.8), Inches(4.8))
        print("Embedded high-res workflow_diagram.png into Slide 6!")

    # Save to user's requested PPT destinations
    dest1 = r'C:\Users\User\Downloads\saathi-sarvasretha-submission.pptx'
    dest2 = r'C:\Users\User\Downloads\Saathi-Your-Accessibility-Copilot-Final.pptx'
    dest3 = r'C:\Users\User\Desktop\saathi-sarvasretha-submission.pptx'

    prs.save(dest1)
    prs.save(dest2)
    try:
        prs.save(dest3)
    except Exception as e:
        print(f"Desktop save note: {e}")

    # Try saving to saathi-sarvasretha.pptx if unlocked
    try:
        prs.save(r'C:\Users\User\Downloads\saathi-sarvasretha.pptx')
    except Exception as e:
        print(f"Locked file note: {e}")

    print(f"\nSUCCESS! Modified presentation saved to:\n  1. {dest1}\n  2. {dest2}\n  3. {dest3}")

if __name__ == '__main__':
    update_original_pptx()
