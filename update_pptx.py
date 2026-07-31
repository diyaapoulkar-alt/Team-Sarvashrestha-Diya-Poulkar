import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_updated_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333) # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # Color Palette
    BG_DARK = RGBColor(15, 23, 42)      # #0f172a
    CARD_BG = RGBColor(30, 41, 59)      # #1e293b
    ACCENT_CYAN = RGBColor(56, 189, 248) # #38bdf8
    ACCENT_PINK = RGBColor(236, 72, 153) # #ec4899
    ACCENT_GREEN = RGBColor(16, 185, 129)# #10b981
    TEXT_LIGHT = RGBColor(248, 250, 252)# #f8fafc
    TEXT_MUTED = RGBColor(203, 213, 225)# #cbd5e1

    def add_bg(slide):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_DARK
        shape.line.fill.background()

    def add_header(slide, title, category="SAATHI ACCESSIBILITY COPILOT | TEAM SARVASHRESTHA"):
        # Header category
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = category
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN

        # Main Slide Title
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = title
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_LIGHT

    # --- SLIDE 1: Title Slide ---
    slide1 = prs.slides.add_slide(blank_layout)
    add_bg(slide1)

    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(3.5))
    tf1 = title_box.text_frame
    
    p = tf1.paragraphs[0]
    p.text = "Saathi — Your Multimodal AI Accessibility Copilot"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    p_sub = tf1.add_paragraph()
    p_sub.text = "Unified Learning Platform with Real-Time Vision AI, Zero-Latency Subtitles, Cognitive Simplifier & AI Lady Tutor"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = TEXT_MUTED
    p_sub.space_before = Pt(14)

    p_meta = tf1.add_paragraph()
    p_meta.text = "PROPOSED BY: TEAM SARVASHRESTHA | TRACK 2: AI AGENT / ACCESSIBILITY COPILOT (SOFC 2.0)"
    p_meta.font.size = Pt(14)
    p_meta.font.bold = True
    p_meta.font.color.rgb = ACCENT_PINK
    p_meta.space_before = Pt(28)

    # --- SLIDE 2: Problem & Solution ---
    slide2 = prs.slides.add_slide(blank_layout)
    add_bg(slide2)
    add_header(slide2, "The Problem & Our Solution")

    # Left Box: Problem
    box_p = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    box_p.fill.solid()
    box_p.fill.fore_color.rgb = CARD_BG
    box_p.line.color.rgb = ACCENT_PINK
    tf_p = box_p.text_frame
    tf_p.word_wrap = True
    p = tf_p.paragraphs[0]
    p.text = "❌ THE PROBLEM"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = ACCENT_PINK
    bullets_p = [
        "Students with disabilities juggle 5+ disconnected tools (screen reader, translator, captioner, PDF simplified reader).",
        "No tool understands the student's unique learning profile across senses.",
        "Existing apps suffer high latency, complex UI setups, or costly subscriptions.",
        "Campuses lack a single unified accessibility layer for inclusive classrooms."
    ]
    for b in bullets_p:
        p_b = tf_p.add_paragraph()
        p_b.text = f"• {b}"
        p_b.font.size = Pt(14); p_b.font.color.rgb = TEXT_MUTED; p_b.space_before = Pt(10)

    # Right Box: Solution
    box_s = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
    box_s.fill.solid()
    box_s.fill.fore_color.rgb = CARD_BG
    box_s.line.color.rgb = ACCENT_GREEN
    tf_s = box_s.text_frame
    tf_s.word_wrap = True
    p = tf_s.paragraphs[0]
    p.text = "✨ THE SAATHI SOLUTION"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = ACCENT_GREEN
    bullets_s = [
        "Single PWA powered by Personalization-Aware Modality Router (PAMR).",
        "Set your accessibility profile once — Visual, Hearing, Cognitive, or Math.",
        "Automatic routing to Vision OCR, 0ms Live Speech Subtitles & 5th-Grade Simplifier.",
        "Includes Saathi AI Studio with Chibi AI Lady Tutor for interactive campus study."
    ]
    for b in bullets_s:
        p_b = tf_s.add_paragraph()
        p_b.text = f"• {b}"
        p_b.font.size = Pt(14); p_b.font.color.rgb = TEXT_MUTED; p_b.space_before = Pt(10)

    # --- SLIDE 3: Key Features & Assistive Profiles ---
    slide3 = prs.slides.add_slide(blank_layout)
    add_bg(slide3)
    add_header(slide3, "Core Assistive Features & Profiles")

    features = [
        ("👁️ Visual Assist", "Llama-3.2 Vision & OCR for printed notes, environment photos & audio screen reader.", ACCENT_CYAN),
        ("🎤 Hearing Assist", "0ms latency live classroom subtitles with instant Hindi/Marathi translation.", ACCENT_PINK),
        ("🧠 Cognitive Assist", "5th-grade textbook simplifier with OpenDyslexic font & rich markdown highlights.", ACCENT_GREEN),
        ("📐 Math & LaTeX Assist", "Phonetic math formula speech reader & symbol breakdown.", RGBColor(245, 158, 11)),
    ]

    for i, (title, desc, color) in enumerate(features):
        x = Inches(0.8 + (i % 2) * 5.9)
        y = Inches(1.8 + (i // 2) * 2.5)
        box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(2.2))
        box.fill.solid(); box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = color
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = color
        p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(14); p2.font.color.rgb = TEXT_MUTED; p2.space_before = Pt(8)

    # --- SLIDE 4: Architecture Diagram Slide ---
    slide4 = prs.slides.add_slide(blank_layout)
    add_bg(slide4)
    add_header(slide4, "System Architecture: PAMR Framework")
    img_path_arch = os.path.join(os.getcwd(), 'public', 'architecture_diagram.png')
    if os.path.exists(img_path_arch):
        slide4.shapes.add_picture(img_path_arch, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.3))

    # --- SLIDE 5: Workflow Flowchart Slide ---
    slide5 = prs.slides.add_slide(blank_layout)
    add_bg(slide5)
    add_header(slide5, "End-to-End User Workflow & Data Pipeline")
    img_path_wf = os.path.join(os.getcwd(), 'public', 'workflow_diagram.png')
    if os.path.exists(img_path_wf):
        slide5.shapes.add_picture(img_path_wf, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.3))

    # --- SLIDE 6: Multilingual UI Engine ---
    slide6 = prs.slides.add_slide(blank_layout)
    add_bg(slide6)
    add_header(slide6, "Full Site Multilingual Translation Engine")
    
    box_lang = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    box_lang.fill.solid(); box_lang.fill.fore_color.rgb = CARD_BG
    box_lang.line.color.rgb = ACCENT_CYAN
    tf_l = box_lang.text_frame; tf_l.word_wrap = True
    p = tf_l.paragraphs[0]; p.text = "🌐 3 SUPPORTED LANGUAGES: ENGLISH, HINDI (हिंदी), MARATHI (मराठी)"; p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = ACCENT_CYAN
    bullets_l = [
        "Instant Dynamic Language Switcher: Changing language updates the ENTIRE PWA interface in real time.",
        "Navbar, Hero Section, Sidebar Navigation, Tool Headers & Profile Controls are fully translated.",
        "Groq AI Engine generates study summaries in the student's selected native language.",
        "Bridges the language barrier for regional Indian university students."
    ]
    for b in bullets_l:
        p_b = tf_l.add_paragraph(); p_b.text = f"• {b}"; p_b.font.size = Pt(16); p_b.font.color.rgb = TEXT_LIGHT; p_b.space_before = Pt(14)

    # --- SLIDE 7: Saathi AI Studio & Lady Tutor ---
    slide7 = prs.slides.add_slide(blank_layout)
    add_bg(slide7)
    add_header(slide7, "Saathi AI Studio & Chibi Lady Tutor")

    # Left: Details
    box_t = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(7.5), Inches(4.8))
    box_t.fill.solid(); box_t.fill.fore_color.rgb = CARD_BG
    box_t.line.color.rgb = ACCENT_PINK
    tf_t = box_t.text_frame; tf_t.word_wrap = True
    p = tf_t.paragraphs[0]; p.text = "👩‍🏫 INTERACTIVE LADY AI TUTOR COPILOT"; p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = ACCENT_PINK
    bullets_t = [
        "Dedicated Saathi AI Studio accessible directly from the left sidebar.",
        "Cute AI Lady Tutor character with smart glasses & pointer stick for friendly learning.",
        "Interactive 1-Click Chips: Pop Quiz Master, 5th-Grade Story, Academic Joke & Study Motivation.",
        "Hands-Free Continuous Speech Loop with automatic voice readout."
    ]
    for b in bullets_t:
        p_b = tf_t.add_paragraph(); p_b.text = f"• {b}"; p_b.font.size = Pt(15); p_b.font.color.rgb = TEXT_MUTED; p_b.space_before = Pt(12)

    # Right: Tutor Image
    img_tutor = os.path.join(os.getcwd(), 'public', 'lady_tutor.png')
    if os.path.exists(img_tutor):
        slide7.shapes.add_picture(img_tutor, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8))

    # --- SLIDE 8: Tech Stack ---
    slide8 = prs.slides.add_slide(blank_layout)
    add_bg(slide8)
    add_header(slide8, "Updated Technical Stack")

    stack_items = [
        ("Frontend Framework", "React 18 + Vite PWA (Zero Install, Progressive Web App)", ACCENT_CYAN),
        ("AI Inference Engine", "Groq Llama-3.3-70b-versatile & Llama-3.2-11b-vision-preview", ACCENT_PINK),
        ("OCR & Speech APIs", "Web Speech API, SpeechSynthesis & Tesseract.js fallback engine", ACCENT_GREEN),
        ("UI & Design System", "Platinum Silver & Dark Charcoal CSS with Glassmorphism & Lucide Icons", RGBColor(245, 158, 11))
    ]

    for i, (title, desc, color) in enumerate(stack_items):
        x = Inches(0.8 + (i % 2) * 5.9)
        y = Inches(1.8 + (i // 2) * 2.5)
        box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(2.2))
        box.fill.solid(); box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = color
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = color
        p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(14); p2.font.color.rgb = TEXT_MUTED; p2.space_before = Pt(8)

    # --- SLIDE 9: Feasibility Analysis ---
    slide9 = prs.slides.add_slide(blank_layout)
    add_bg(slide9)
    add_header(slide9, "Feasibility & Zero-Cost Architecture")

    box_f = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    box_f.fill.solid(); box_f.fill.fore_color.rgb = CARD_BG
    box_f.line.color.rgb = ACCENT_GREEN
    tf_f = box_f.text_frame; tf_f.word_wrap = True
    p = tf_f.paragraphs[0]; p.text = "⚡ HACKATHON FEASIBILITY & DEPLOYMENT"; p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = ACCENT_GREEN
    bullets_f = [
        "1. 4-Person Team: Roles clearly scoped across AI, Frontend, Accessibility, and Design.",
        "2. Zero-Cost Infrastructure: Groq AI, Vercel hosting & Web Speech API — 100% free tier.",
        "3. Zero Hardware Barrier: Runs on any browser phone/laptop without special hardware.",
        "4. Pre-trained AI Models: Zero custom model training required; instant high accuracy out of the box."
    ]
    for b in bullets_f:
        p_b = tf_f.add_paragraph(); p_b.text = b; p_b.font.size = Pt(15); p_b.font.color.rgb = TEXT_LIGHT; p_b.space_before = Pt(14)

    # --- SLIDE 10: Risk Mitigation ---
    slide10 = prs.slides.add_slide(blank_layout)
    add_bg(slide10)
    add_header(slide10, "Risk Mitigation & Live Demo Readiness")

    box_r = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    box_r.fill.solid(); box_r.fill.fore_color.rgb = CARD_BG
    box_r.line.color.rgb = ACCENT_PINK
    tf_r = box_r.text_frame; tf_r.word_wrap = True
    p = tf_r.paragraphs[0]; p.text = "🛡️ OVERCOMING TECHNICAL CHALLENGES"; p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = ACCENT_PINK
    bullets_r = [
        "Speech Latency: Progressive streaming & client-side SpeechSynthesis voice buffer.",
        "Image / Note Noise: Built-in filter Garbled OCR regex engine cleans illegible symbols.",
        "Guest Authentication: 1-Click Guest Sign In for instant reviewer access without sign-up delay.",
        "Live Deployment: Deployed on Vercel with automatic GitHub sync."
    ]
    for b in bullets_r:
        p_b = tf_r.add_paragraph(); p_b.text = f"• {b}"; p_b.font.size = Pt(15); p_b.font.color.rgb = TEXT_MUTED; p_b.space_before = Pt(14)

    # --- SLIDE 11: Real-World Impact ---
    slide11 = prs.slides.add_slide(blank_layout)
    add_bg(slide11)
    add_header(slide11, "Impact & Campus Benefits")

    box_imp = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    box_imp.fill.solid(); box_imp.fill.fore_color.rgb = CARD_BG
    box_imp.line.color.rgb = ACCENT_CYAN
    tf_imp = box_imp.text_frame; tf_imp.word_wrap = True
    p = tf_imp.paragraphs[0]; p.text = "🌟 REAL-WORLD CAMPUS IMPACT"; p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = ACCENT_CYAN
    bullets_imp = [
        "Replaces 5+ fragmented apps with one unified accessibility PWA.",
        "Empowers visually impaired, deaf/hard-of-hearing, and dyslexic students in university classrooms.",
        "Zero-install PWA eliminates device compatibility barriers for low-income students.",
        "Supports non-native English speakers with live Hindi and Marathi translation."
    ]
    for b in bullets_imp:
        p_b = tf_imp.add_paragraph(); p_b.text = f"• {b}"; p_b.font.size = Pt(15); p_b.font.color.rgb = TEXT_LIGHT; p_b.space_before = Pt(14)

    # --- SLIDE 12: Why Choose Saathi ---
    slide12 = prs.slides.add_slide(blank_layout)
    add_bg(slide12)
    add_header(slide12, "Why Choose Saathi?")

    box_w = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    box_w.fill.solid(); box_w.fill.fore_color.rgb = CARD_BG
    box_w.line.color.rgb = ACCENT_GREEN
    tf_w = box_w.text_frame; tf_w.word_wrap = True
    p = tf_w.paragraphs[0]; p.text = "💡 THE PAMR DIFFERENCE"; p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = ACCENT_GREEN
    bullets_w = [
        "1. All-in-One Platform: Vision, Speech, Text & Math tools in one unified PWA.",
        "2. Automatic AI Routing: PAMR eliminates manual settings — set your profile once.",
        "3. Browser-Based PWA: Zero app store installation needed — works instantly on any phone or laptop."
    ]
    for b in bullets_w:
        p_b = tf_w.add_paragraph(); p_b.text = b; p_b.font.size = Pt(16); p_b.font.color.rgb = TEXT_LIGHT; p_b.space_before = Pt(16)

    # --- SLIDE 13: Academic References ---
    slide13 = prs.slides.add_slide(blank_layout)
    add_bg(slide13)
    add_header(slide13, "Academic References & Research Foundations")

    box_ref = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    box_ref.fill.solid(); box_ref.fill.fore_color.rgb = CARD_BG
    box_ref.line.color.rgb = RGBColor(245, 158, 11)
    tf_ref = box_ref.text_frame; tf_ref.word_wrap = True
    p = tf_ref.paragraphs[0]; p.text = "📚 RESEARCH PAPERS & FOUNDATIONS"; p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = RGBColor(245, 158, 11)
    refs = [
        "Manduchi, R., & Coughlan, J. (2012). 'The Role of Computer Vision in Accessible Technology.' Proceedings of the IEEE, 100(5), 1680–1698.",
        "Karamolegkou, A., et al. (2025). 'Evaluating multimodal language models as visual assistants for visually impaired users.' ACL 2025 / arXiv.",
        "Microsoft Research — 'Seeing AI: Helping People See the World' (real-world assistive framework reference)."
    ]
    for r in refs:
        p_b = tf_ref.add_paragraph(); p_b.text = f"• {r}"; p_b.font.size = Pt(15); p_b.font.color.rgb = TEXT_MUTED; p_b.space_before = Pt(14)

    # --- SLIDE 14: Future Roadmap ---
    slide14 = prs.slides.add_slide(blank_layout)
    add_bg(slide14)
    add_header(slide14, "Future Roadmap & Vision")

    box_m = slide14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    box_m.fill.solid(); box_m.fill.fore_color.rgb = CARD_BG
    box_m.line.color.rgb = ACCENT_CYAN
    tf_m = box_m.text_frame; tf_m.word_wrap = True
    p = tf_m.paragraphs[0]; p.text = "🚀 SCALABILITY & VISION"; p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = ACCENT_CYAN
    roadmap = [
        "Short-Term: Add 10+ regional Indian languages & offline web worker OCR mode.",
        "Mid-Term: Smart Classroom LMS integration & AI campus navigation.",
        "Long-Term: Wearable glass camera & smart mic integration for hands-free university lectures.",
        "Mission: 'To build an inclusive digital learning ecosystem where every student accesses education independently.'"
    ]
    for r in roadmap:
        p_b = tf_m.add_paragraph(); p_b.text = f"• {r}"; p_b.font.size = Pt(15); p_b.font.color.rgb = TEXT_LIGHT; p_b.space_before = Pt(14)

    # --- SLIDE 15: Team Sarvashrestha ---
    slide15 = prs.slides.add_slide(blank_layout)
    add_bg(slide15)
    add_header(slide15, "Team Sarvashrestha")

    box_team = slide15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    box_team.fill.solid(); box_team.fill.fore_color.rgb = CARD_BG
    box_team.line.color.rgb = ACCENT_PINK
    tf_t = box_team.text_frame; tf_t.word_wrap = True
    p = tf_t.paragraphs[0]; p.text = "👥 TEAM MEMBERS & CONTRIBUTIONS"; p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = ACCENT_PINK
    members = [
        "Diya Poulkar & Team Sarvashrestha",
        "Track 2: AI Agent & Accessibility Copilot",
        "Repository: https://github.com/diyaapoulkar-alt/Team-Sarvashrestha-Diya-Poulkar.git",
        "Live Deployment: https://team-sarvashrestha-diya-poulkar.vercel.app/"
    ]
    for m in members:
        p_b = tf_t.add_paragraph(); p_b.text = f"• {m}"; p_b.font.size = Pt(16); p_b.font.color.rgb = TEXT_LIGHT; p_b.space_before = Pt(16)

    # --- SLIDE 16: Thank You Slide ---
    slide16 = prs.slides.add_slide(blank_layout)
    add_bg(slide16)

    t_box = slide16.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
    tf16 = t_box.text_frame
    p = tf16.paragraphs[0]; p.text = "THANK YOU!"; p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = ACCENT_CYAN
    p2 = tf16.add_paragraph(); p2.text = "Saathi — Universal Learning with Multimodal AI Accessibility"; p2.font.size = Pt(20); p2.font.color.rgb = TEXT_LIGHT; p2.space_before = Pt(16)
    p3 = tf16.add_paragraph(); p3.text = "Live Demo: https://team-sarvashrestha-diya-poulkar.vercel.app/"; p3.font.size = Pt(16); p3.font.color.rgb = ACCENT_GREEN; p3.space_before = Pt(16)

    # Save outputs
    out_dl = r'C:\Users\User\Downloads\saathi-sarvasretha.pptx'
    out_pub = os.path.join(os.getcwd(), 'public', 'saathi-sarvasretha.pptx')
    out_ast = os.path.join(os.getcwd(), 'src', 'assets', 'saathi-sarvasretha.pptx')

    prs.save(out_dl)
    prs.save(out_pub)
    prs.save(out_ast)
    print("Successfully updated PowerPoint PPTX presentation across Downloads and Public assets!")

if __name__ == '__main__':
    create_updated_presentation()
